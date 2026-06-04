"""Designer — import PPTX and DOCX files as designer pages.

Converts uploaded documents into DesignerPage objects with HTML content:
  - PPTX: each slide → one page with positioned text boxes
  - DOCX: paragraphs grouped into pages (split at headings or every ~800 words)
"""

from __future__ import annotations

import base64
import io
import logging
import re
from typing import Optional

from row_bot.designer.state import DesignerPage

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════
# PPTX IMPORT
# ═══════════════════════════════════════════════════════════════════════

def import_pptx(file_bytes: bytes) -> list[DesignerPage]:
    """Convert a PPTX file into a list of DesignerPage objects.

    Each slide becomes one page with text boxes positioned via absolute CSS.
    Images are embedded as base64.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError("python-pptx is required for PPTX import. Run: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    slide_w = prs.slide_width or Emu(12192000)  # default 10" in EMU
    slide_h = prs.slide_height or Emu(6858000)   # default 7.5"
    w_px = int(slide_w / 914400 * 96)  # EMU → inches → pixels (96 DPI)
    h_px = int(slide_h / 914400 * 96)

    pages = []
    for i, slide in enumerate(prs.slides):
        elements_html = []

        for shape in slide.shapes:
            left = int((shape.left or 0) / 914400 * 96)
            top = int((shape.top or 0) / 914400 * 96)
            width = int((shape.width or 0) / 914400 * 96)
            height = int((shape.height or 0) / 914400 * 96)

            if shape.has_text_frame:
                paragraphs_html = []
                for para in shape.text_frame.paragraphs:
                    runs_html = []
                    for run in para.runs:
                        style_parts = []
                        if run.font.bold:
                            style_parts.append("font-weight:bold")
                        if run.font.italic:
                            style_parts.append("font-style:italic")
                        if run.font.size:
                            size_pt = run.font.size.pt
                            style_parts.append(f"font-size:{size_pt}pt")
                        if run.font.color and run.font.color.type != None and run.font.color.rgb:
                            style_parts.append(f"color:#{run.font.color.rgb}")
                        text = _escape_html(run.text)
                        if style_parts:
                            runs_html.append(f'<span style="{";".join(style_parts)}">{text}</span>')
                        else:
                            runs_html.append(text)

                    align = "left"
                    if para.alignment:
                        align_map = {
                            PP_ALIGN.CENTER: "center",
                            PP_ALIGN.RIGHT: "right",
                            PP_ALIGN.JUSTIFY: "justify",
                        }
                        align = align_map.get(para.alignment, "left")
                    p_content = "".join(runs_html)
                    if p_content.strip():
                        paragraphs_html.append(f'<p style="text-align:{align};margin:0 0 4px 0;">{p_content}</p>')

                if paragraphs_html:
                    inner = "\n".join(paragraphs_html)
                    elements_html.append(
                        f'<div style="position:absolute;left:{left}px;top:{top}px;'
                        f'width:{width}px;height:{height}px;overflow:hidden;">\n'
                        f'{inner}\n</div>'
                    )

            elif shape.shape_type == 13:  # Picture
                try:
                    img_bytes = shape.image.blob
                    content_type = shape.image.content_type or "image/png"
                    b64 = base64.b64encode(img_bytes).decode("ascii")
                    elements_html.append(
                        f'<img src="data:{content_type};base64,{b64}" '
                        f'style="position:absolute;left:{left}px;top:{top}px;'
                        f'width:{width}px;height:{height}px;object-fit:contain;" />'
                    )
                except Exception:
                    logger.debug("Skipping image in slide %d", i + 1)

        # Extract speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text or ""

        # Build page title from first heading-like text
        title = f"Slide {i + 1}"
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()[:60]
                break

        body = "\n".join(elements_html) if elements_html else '<p style="text-align:center;padding-top:40%;color:#888;">Empty slide</p>'
        page_html = (
            f'<!DOCTYPE html><html><head><meta charset="utf-8">'
            f'<style>* {{margin:0;padding:0;box-sizing:border-box;}} '
            f'body {{width:{w_px}px;height:{h_px}px;position:relative;'
            f'font-family:var(--body-font,"Inter",sans-serif);'
            f'background:var(--bg,#fff);color:var(--text,#1a1a1a);}}</style>'
            f'</head><body>\n{body}\n</body></html>'
        )

        pages.append(DesignerPage(html=page_html, title=title, notes=notes))

    if not pages:
        pages.append(DesignerPage(title="Imported (empty)", html="<html><body><p>No slides found.</p></body></html>"))

    logger.info("Imported %d slides from PPTX", len(pages))
    return pages


# ═══════════════════════════════════════════════════════════════════════
# DOCX IMPORT
# ═══════════════════════════════════════════════════════════════════════

def import_docx(file_bytes: bytes) -> list[DesignerPage]:
    """Convert a DOCX file into a list of DesignerPage objects.

    Splits at Heading 1 / Heading 2 boundaries, or every ~800 words.
    """
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        raise RuntimeError("python-docx is required for DOCX import. Run: pip install python-docx")

    doc = Document(io.BytesIO(file_bytes))

    # Group paragraphs into sections (split at H1/H2)
    sections: list[dict] = []
    current: dict = {"title": "Page 1", "paragraphs": []}
    word_count = 0

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower()
        text = para.text.strip()
        if not text:
            current["paragraphs"].append("<br/>")
            continue

        # Split on headings
        if style_name.startswith("heading 1") or style_name.startswith("heading 2"):
            if current["paragraphs"]:
                sections.append(current)
            current = {"title": text[:60], "paragraphs": []}
            word_count = 0
            tag = "h1" if "1" in style_name else "h2"
            current["paragraphs"].append(f"<{tag}>{_escape_html(text)}</{tag}>")
            continue

        # Split on word count threshold
        word_count += len(text.split())
        if word_count > 800 and current["paragraphs"]:
            sections.append(current)
            current = {"title": f"Page {len(sections) + 1}", "paragraphs": []}
            word_count = len(text.split())

        # Convert paragraph
        html_content = _docx_para_to_html(para, style_name)
        current["paragraphs"].append(html_content)

    if current["paragraphs"]:
        sections.append(current)

    # Also extract tables
    for table in doc.tables:
        rows_html = []
        for row in table.rows:
            cells = "".join(f"<td style='padding:8px;border:1px solid #ddd;'>{_escape_html(c.text)}</td>" for c in row.cells)
            rows_html.append(f"<tr>{cells}</tr>")
        table_html = f"<table style='border-collapse:collapse;width:100%;margin:16px 0;'>{''.join(rows_html)}</table>"
        if sections:
            sections[-1]["paragraphs"].append(table_html)
        else:
            sections.append({"title": "Table", "paragraphs": [table_html]})

    # Convert sections to pages
    pages = []
    for sec in sections:
        body = "\n".join(sec["paragraphs"])
        page_html = (
            '<!DOCTYPE html><html><head><meta charset="utf-8">'
            '<style>* {margin:0;padding:0;box-sizing:border-box;} '
            'body {padding:60px 80px;font-family:var(--body-font,"Inter",sans-serif);'
            'background:var(--bg,#fff);color:var(--text,#1a1a1a);line-height:1.7;} '
            'h1 {font-family:var(--heading-font,"Inter",sans-serif);font-size:36px;margin:0 0 24px 0;} '
            'h2 {font-family:var(--heading-font,"Inter",sans-serif);font-size:28px;margin:24px 0 16px 0;} '
            'h3 {font-family:var(--heading-font,"Inter",sans-serif);font-size:22px;margin:20px 0 12px 0;} '
            'p {margin:0 0 12px 0;font-size:16px;} '
            'ul,ol {margin:0 0 12px 24px;} '
            'li {margin:0 0 4px 0;}</style>'
            f'</head><body>\n{body}\n</body></html>'
        )
        pages.append(DesignerPage(html=page_html, title=sec["title"]))

    if not pages:
        pages.append(DesignerPage(title="Imported (empty)", html="<html><body><p>No content found.</p></body></html>"))

    logger.info("Imported %d pages from DOCX", len(pages))
    return pages


def _docx_para_to_html(para, style_name: str) -> str:
    """Convert a single docx paragraph to HTML."""
    try:
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        WD_ALIGN_PARAGRAPH = None

    runs_html = []
    for run in para.runs:
        text = _escape_html(run.text)
        if not text:
            continue
        styles = []
        if run.bold:
            text = f"<strong>{text}</strong>"
        if run.italic:
            text = f"<em>{text}</em>"
        if run.underline:
            text = f"<u>{text}</u>"
        runs_html.append(text)

    content = "".join(runs_html)
    if not content.strip():
        return "<br/>"

    if "heading 3" in style_name:
        return f"<h3>{content}</h3>"
    if "heading 4" in style_name:
        return f"<h4>{content}</h4>"
    if "list" in style_name:
        return f"<li>{content}</li>"

    align = ""
    if WD_ALIGN_PARAGRAPH and para.alignment:
        align_map = {
            WD_ALIGN_PARAGRAPH.CENTER: "text-align:center;",
            WD_ALIGN_PARAGRAPH.RIGHT: "text-align:right;",
            WD_ALIGN_PARAGRAPH.JUSTIFY: "text-align:justify;",
        }
        align = align_map.get(para.alignment, "")

    if align:
        return f'<p style="{align}">{content}</p>'
    return f"<p>{content}</p>"


def _escape_html(text: str) -> str:
    """Escape HTML special characters."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )
