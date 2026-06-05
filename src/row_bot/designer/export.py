"""Designer — export pipeline (PDF, PPTX, HTML, PNG)."""

from __future__ import annotations

import io
import logging
import math
import os
import pathlib
import re
import zipfile
from typing import Any, Optional

from row_bot.designer.preview import render_page_html
from row_bot.designer.state import DesignerProject

logger = logging.getLogger(__name__)

_WORKSPACE = pathlib.Path(
    os.environ.get("ROW_BOT_WORKSPACE", pathlib.Path.home() / "Documents" / "Row-Bot")
)


class ExportedBytes(bytes):
    """Bytes payload annotated with the actual path written to disk."""

    saved_path: pathlib.Path

    def __new__(cls, payload: bytes, saved_path: pathlib.Path):
        obj = super().__new__(cls, payload)
        obj.saved_path = pathlib.Path(saved_path)
        return obj

_CSS_NUMBER_RE = re.compile(r"-?\d+(?:\.\d+)?")
_RENDERED_DOM_EXPORT_SCRIPT = r"""
() => {
    const normalizeText = (value) => (value || '').replace(/\u00a0/g, ' ').replace(/\s+/g, ' ').trim();
    const parseNumber = (value) => {
        const parsed = Number.parseFloat(value || '0');
        return Number.isFinite(parsed) ? parsed : 0;
    };
    const parseAlpha = (value) => {
        if (!value || value === 'transparent') {
            return 0;
        }
        const text = String(value).trim().toLowerCase();
        if (text.startsWith('#')) {
            return 1;
        }
        const parts = text.match(/[\d.]+/g) || [];
        if (parts.length < 3) {
            return 0;
        }
        return parts.length >= 4 ? Math.max(0, Math.min(1, Number.parseFloat(parts[3]) || 0)) : 1;
    };
    const transformText = (value, transform) => {
        if (!value) {
            return '';
        }
        switch ((transform || '').toLowerCase()) {
            case 'uppercase':
                return value.toUpperCase();
            case 'lowercase':
                return value.toLowerCase();
            case 'capitalize':
                return value.replace(/\b(\w)/g, (match) => match.toUpperCase());
            default:
                return value;
        }
    };
        const INLINE_TAGS = new Set(['span', 'em', 'strong', 'b', 'i', 'u', 'small', 'mark', 'sup', 'sub', 'code', 'a', 'label', 'time', 'abbr', 'cite', 'kbd', 'q', 's', 'var']);
    const ownText = (element) => Array.from(element.childNodes)
        .filter((node) => node.nodeType === Node.TEXT_NODE)
        .map((node) => node.textContent || '')
        .join(' ');
        const isInlineLike = (tag, style) => {
            const display = (style?.display || '').toLowerCase();
            // Any display variant that starts with "inline" (inline,
            // inline-block, inline-flex, inline-grid, inline-table) is
            // flowed with surrounding text and must NOT trigger the
            // "this is a block child" branch in the dedup logic.
            // Treating inline-block as block was the primary cause of
            // double-rendered headings and eyebrow labels in the
            // editable PPTX export.
            return INLINE_TAGS.has(tag) || display === 'contents' || display.startsWith('inline');
        };
    const normalizedVisibleText = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        return normalizeText(transformText(element.innerText || '', resolvedStyle.textTransform));
    };
    const ownVisibleText = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        return normalizeText(transformText(ownText(element), resolvedStyle.textTransform));
    };
    const textBearingChildElements = (element) => Array.from(element.children).filter((child) => {
        if (!(child instanceof HTMLElement)) {
            return false;
        }
        return !!normalizedVisibleText(child);
    });
    const hasBlockTextDescendant = (element) => Array.from(element.querySelectorAll('*')).some((child) => {
        if (!(child instanceof HTMLElement)) {
            return false;
        }
        const childText = normalizedVisibleText(child);
        if (!childText) {
            return false;
        }
        const childStyle = getComputedStyle(child);
        return !isInlineLike(child.tagName.toLowerCase(), childStyle);
    });
    const ownsLeafTextCluster = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        if (hasBlockTextDescendant(element)) {
            return false;
        }
        return !!ownVisibleText(element, resolvedStyle) || textBearingChildElements(element).length <= 1;
    };
        const exportedText = (element, tag, style) => {
        const visibleText = normalizedVisibleText(element, style);
            if (!visibleText) {
                return '';
            }

        const ownTextValue = ownVisibleText(element, style);
        const blockTextDescendant = hasBlockTextDescendant(element);

            if (!isInlineLike(tag, style)) {
            if (!blockTextDescendant) {
                return ownsLeafTextCluster(element, style) ? visibleText : '';
            }
            return ownTextValue;
            }

            const parent = element.parentElement;
            if (!parent) {
                return visibleText;
            }

            const parentStyle = getComputedStyle(parent);
        const parentText = normalizedVisibleText(parent, parentStyle);
        return parentText && ownsLeafTextCluster(parent, parentStyle) ? '' : visibleText;
        };
    const isVisible = (style, rect) => {
        if (!style || style.display === 'none' || style.visibility === 'hidden') {
            return false;
        }
        if (parseFloat(style.opacity || '1') <= 0.01) {
            return false;
        }
        if (rect.width < 4 || rect.height < 4) {
            return false;
        }
        if (rect.right <= 0 || rect.bottom <= 0) {
            return false;
        }
        return true;
    };

    const items = [];
    let order = 0;
    let screenshotCounter = 0;

    for (const element of Array.from(document.body.querySelectorAll('*'))) {
        const rect = element.getBoundingClientRect();
        const style = getComputedStyle(element);
        if (!isVisible(style, rect)) {
            continue;
        }

        const tag = element.tagName.toLowerCase();
        const text = exportedText(element, tag, style);
        const backgroundColor = style.backgroundColor || '';
        const backgroundImage = style.backgroundImage || 'none';
        const borderTopWidth = parseNumber(style.borderTopWidth);
        const borderRightWidth = parseNumber(style.borderRightWidth);
        const borderBottomWidth = parseNumber(style.borderBottomWidth);
        const borderLeftWidth = parseNumber(style.borderLeftWidth);
        const borderWidth = Math.max(borderTopWidth, borderRightWidth, borderBottomWidth, borderLeftWidth);
        const borderColor = style.borderTopColor || style.borderColor || '';
        const borderStyle = style.borderTopStyle || style.borderStyle || 'none';
        const fillAlpha = parseAlpha(backgroundColor);
        const hasFill = fillAlpha > 0.02;
        const hasBorder = borderStyle !== 'none' && borderWidth > 0.1 && parseAlpha(borderColor) > 0.02;
        const hasBackgroundImage = backgroundImage !== 'none';
        const boxShadow = style.boxShadow || 'none';
        const hasBoxShadow = boxShadow !== 'none' && boxShadow.trim() !== '';
        const hasVisual = hasFill || hasBorder || hasBackgroundImage || hasBoxShadow;
        const zIndexRaw = Number.parseInt(style.zIndex || '0', 10);
        const zIndex = Number.isFinite(zIndexRaw) ? zIndexRaw : 0;

        const base = {
            order: order++,
            tag,
            x: rect.left,
            y: rect.top,
            width: rect.width,
            height: rect.height,
            zIndex,
            backgroundColor,
            backgroundImage,
            borderColor,
            borderWidth,
            borderTopWidth,
            borderRightWidth,
            borderBottomWidth,
            borderLeftWidth,
            borderRadius: style.borderRadius || '0px',
            borderTopLeftRadius: style.borderTopLeftRadius || '0px',
            borderTopRightRadius: style.borderTopRightRadius || '0px',
            borderBottomLeftRadius: style.borderBottomLeftRadius || '0px',
            borderBottomRightRadius: style.borderBottomRightRadius || '0px',
            boxShadow,
            fillAlpha,
            color: style.color || '',
            fontFamily: style.fontFamily || '',
            fontSize: style.fontSize || '',
            fontStyle: style.fontStyle || 'normal',
            fontWeight: style.fontWeight || '400',
            lineHeight: style.lineHeight || '',
            textAlign: style.textAlign || 'left',
            opacity: parseFloat(style.opacity || '1') || 1,
            paddingTop: style.paddingTop || '0px',
            paddingRight: style.paddingRight || '0px',
            paddingBottom: style.paddingBottom || '0px',
            paddingLeft: style.paddingLeft || '0px',
            display: style.display || 'block',
            alignItems: style.alignItems || 'stretch',
            justifyContent: style.justifyContent || 'flex-start',
            whiteSpace: style.whiteSpace || 'normal',
        };

        // Raster kinds: <img> bitmaps and inline <svg> roots.  Both
        // are captured as a per-element screenshot (bounded — neither
        // has a visible HTML descendant we'd also emit, so there is
        // no ancestor/descendant duplication class).  Nested <svg>
        // elements are skipped so we only screenshot the outermost.
        if (tag === 'img' || (tag === 'svg' && !element.parentElement?.closest('svg'))) {
            const screenshotId = `pptx-node-${screenshotCounter++}`;
            element.setAttribute('data-row-bot-pptx-export-id', screenshotId);
            const extras = {};
            if (tag === 'img') {
                extras.src = element.src || '';
                extras.objectFit = style.objectFit || 'fill';
                extras.objectPosition = style.objectPosition || '50% 50%';
            } else {
                extras.svgOuterHTML = element.outerHTML || '';
            }
            items.push({
                ...base,
                kind: 'raster',
                screenshotId,
                ...extras,
            });
            continue;
        }

        // Skip descendants of SVG — they are captured in the parent
        // SVG screenshot.
        if (element.closest('svg') && tag !== 'svg') {
            continue;
        }

        // Shape: any element with a fill, border, gradient, or shadow.
        // We no longer try to decide between "native shape" and
        // "raster snapshot" — every element that has a visual becomes
        // a native PPTX auto-shape.  Gradients and shadows are emitted
        // via lxml on the Python side.  The result is fully editable
        // in PowerPoint, at the cost of some visual drift on complex
        // CSS effects (acceptable tradeoff for Editable mode; users
        // who want pixel-perfect output pick High-Fidelity mode).
        if (hasVisual && rect.width >= 6 && rect.height >= 6) {
            items.push({
                ...base,
                kind: 'shape',
            });
        }

        if (text) {
            items.push({
                ...base,
                kind: 'text',
                text: tag === 'li' && !text.startsWith('•') ? `• ${text}` : text,
            });
        }
    }

    const bodyStyle = getComputedStyle(document.body);
    const rootStyle = getComputedStyle(document.documentElement);
    const bodyBackground = parseAlpha(bodyStyle.backgroundColor) > 0.02
        ? bodyStyle.backgroundColor
        : rootStyle.backgroundColor;
    const bodyBackgroundImage = bodyStyle.backgroundImage && bodyStyle.backgroundImage !== 'none'
        ? bodyStyle.backgroundImage
        : (rootStyle.backgroundImage && rootStyle.backgroundImage !== 'none' ? rootStyle.backgroundImage : '');

    return {
        backgroundColor: bodyBackground || '',
        backgroundImage: bodyBackgroundImage,
        items,
    };
}
"""


def get_export_workspace() -> pathlib.Path:
    """Return the default export workspace and ensure it exists."""
    _WORKSPACE.mkdir(parents=True, exist_ok=True)
    return _WORKSPACE


def _sanitize_name(text: str, max_len: int = 60) -> str:
    safe = "".join(c if c.isalnum() or c in " -_" else "_" for c in (text or ""))
    safe = safe.strip()[:max_len]
    return safe or "Designer Export"


def _resolve_directory(directory: pathlib.Path | str | os.PathLike | None) -> pathlib.Path:
    if directory is None:
        return get_export_workspace()
    resolved = pathlib.Path(directory)
    resolved.mkdir(parents=True, exist_ok=True)
    return resolved


def _parse_page_range(pages_str: Optional[str], total: int) -> list[int]:
    """Parse a page range string like '1-3' or '1,3,5' into 0-based indices."""
    if not pages_str or pages_str.lower() == "all":
        return list(range(total))
    indices = set()
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = max(0, int(a.strip()) - 1)
            end = min(total, int(b.strip()))
            indices.update(range(start, end))
        else:
            idx = int(part.strip()) - 1
            if 0 <= idx < total:
                indices.add(idx)
    return sorted(indices) if indices else list(range(total))


def _selected_pages(project: DesignerProject, pages: Optional[str]) -> list[tuple[int, object]]:
    indices = _parse_page_range(pages, len(project.pages))
    return [(i, project.pages[i]) for i in indices if i < len(project.pages)]


def describe_export_destination(
    project: DesignerProject,
    format: str,
    pages: Optional[str] = None,
    mode: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> pathlib.Path:
    """Return the exact path an export operation will write to."""
    fmt = (format or "").lower().strip()
    out_dir = _resolve_directory(directory)
    safe_name = _sanitize_name(project.name)
    page_count = len(_parse_page_range(pages, len(project.pages)))
    if fmt == "html":
        filename = f"{safe_name}.html"
    elif fmt == "pdf":
        filename = f"{safe_name}.pdf"
    elif fmt == "png":
        filename = f"{safe_name}.png" if page_count <= 1 else f"{safe_name}_pages.zip"
    elif fmt == "pptx":
        filename = f"{safe_name}_editable.pptx" if (mode or "").lower() == "structured" else f"{safe_name}.pptx"
    else:
        raise ValueError(f"Unsupported export format: {format}")
    return out_dir / filename


def _next_available_export_path(path: pathlib.Path) -> pathlib.Path:
    candidate = path
    counter = 2
    while candidate.exists():
        candidate = path.with_name(f"{path.stem} ({counter}){path.suffix}")
        counter += 1
    return candidate


def _permission_denied_message(path: pathlib.Path, label: str) -> str:
    return (
        f"{label} could not be saved to {path} because that file is open or locked. "
        "Close it in PowerPoint or Explorer preview and try again."
    )


def _save_bytes(path: pathlib.Path, data: bytes, label: str) -> bytes:
    path.parent.mkdir(parents=True, exist_ok=True)
    saved_path = path
    try:
        path.write_bytes(data)
    except PermissionError as exc:
        if path.exists():
            fallback_path = _next_available_export_path(path)
            try:
                fallback_path.write_bytes(data)
            except PermissionError as fallback_exc:
                raise PermissionError(_permission_denied_message(path, label)) from fallback_exc
            saved_path = fallback_path
            logger.warning("%s target locked; saved to %s instead of %s", label, saved_path, path)
        else:
            raise PermissionError(_permission_denied_message(path, label)) from exc

    logger.info("Exported %s to %s (%d bytes)", label, saved_path, len(data))
    return ExportedBytes(data, saved_path)


def _escape_html(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _escape_attr(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace('"', "&quot;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def build_html_export(project: DesignerProject, pages: Optional[str] = None) -> bytes:
    """Bundle all or selected pages into a single self-contained HTML file."""
    selected = _selected_pages(project, pages)

    from row_bot.designer.fonts import get_font_css_embedded

    brand = project.brand
    font_families = list(dict.fromkeys([brand.heading_font, brand.body_font] if brand else []))
    embedded_font_css = "\n".join(get_font_css_embedded(f) for f in font_families)

    sections = []
    for i, page in selected:
        html = render_page_html(project, page.html, page_index=i)
        sections.append(
            f'<section id="page-{i}" style="margin-bottom:40px; page-break-after:always;">\n'
            f'<h2 style="font-family:sans-serif;font-size:14px;color:#888;margin-bottom:8px;">'
            f'Page {i + 1}: {page.title}</h2>\n'
            f'<div style="border:1px solid #333;border-radius:8px;overflow:hidden;">\n'
            f'<iframe srcdoc="{_escape_attr(html)}" '
            f'style="width:{project.canvas_width}px;height:{project.canvas_height}px;border:none;" '
            f'sandbox="allow-same-origin"></iframe>\n'
            f'</div>\n</section>'
        )

    indices = _parse_page_range(pages, len(project.pages))
    nav_links = " ".join(
        f'<a href="#page-{i}" style="color:#2563EB;margin-right:12px;">'
        f'{project.pages[i].title}</a>'
        for i in indices if i < len(project.pages)
    )

    full = (
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<title>{_escape_html(project.name)}</title>"
        f"<style>{embedded_font_css}\n"
        f"body{{background:#111;color:#fff;font-family:sans-serif;padding:20px;}}"
        f"a{{text-decoration:none;}}</style></head><body>"
        f"<h1>{_escape_html(project.name)}</h1>"
        f"<nav style='margin-bottom:20px;'>{nav_links}</nav>"
        + "\n".join(sections)
        + "</body></html>"
    )
    return full.encode("utf-8")


def export_html(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    data = build_html_export(project, pages)
    out_path = describe_export_destination(project, "html", pages, directory=directory)
    return _save_bytes(out_path, data, "HTML")


def export_pdf(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render each page as a PDF using Playwright and merge the result."""
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError(
            "Playwright is required for PDF export. Run: pip install playwright && playwright install chromium"
        ) from exc

    pdf_pages = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height}
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            pdf_pages.append(
                pg.pdf(
                    width=f"{project.canvas_width}px",
                    height=f"{project.canvas_height}px",
                    print_background=True,
                )
            )
            ctx.close()
        browser.close()

    if len(pdf_pages) == 1:
        merged = pdf_pages[0]
    else:
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            logger.warning("pypdf not installed — returning first page only for multi-page PDF export")
            merged = pdf_pages[0]
        else:
            writer = PdfWriter()
            for pdf_data in pdf_pages:
                reader = PdfReader(io.BytesIO(pdf_data))
                for page in reader.pages:
                    writer.add_page(page)
            buf = io.BytesIO()
            writer.write(buf)
            merged = buf.getvalue()

    out_path = describe_export_destination(project, "pdf", pages, directory=directory)
    return _save_bytes(out_path, merged, "PDF")


def _render_png_screenshots(project: DesignerProject, pages: Optional[str] = None) -> list[tuple[str, bytes]]:
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError("Playwright is required for PNG export.") from exc

    screenshots: list[tuple[str, bytes]] = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height},
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            png_bytes = pg.screenshot(full_page=False, type="png")
            safe_title = _sanitize_name(page.title, max_len=40)
            screenshots.append((f"page_{i + 1}_{safe_title}.png", png_bytes))
            ctx.close()
        browser.close()
    return screenshots


def export_png_files(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> list[pathlib.Path]:
    """Write selected pages as individual PNG files and return the paths."""
    screenshots = _render_png_screenshots(project, pages)
    out_dir = _resolve_directory(directory)
    paths: list[pathlib.Path] = []
    for filename, png_bytes in screenshots:
        path = out_dir / filename
        path.write_bytes(png_bytes)
        paths.append(path)
    logger.info("Exported %d individual PNG files to %s", len(paths), out_dir)
    return paths


def export_png(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Screenshot selected pages as PNG. Returns one PNG or a ZIP for multiple pages."""
    screenshots = _render_png_screenshots(project, pages)
    out_path = describe_export_destination(project, "png", pages, directory=directory)

    if len(screenshots) == 1:
        return _save_bytes(out_path, screenshots[0][1], "PNG")

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as archive:
        for filename, png_bytes in screenshots:
            archive.writestr(filename, png_bytes)
    return _save_bytes(out_path, buf.getvalue(), "PNG ZIP")


def export_pptx_screenshot(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render each page as an image via Playwright and embed in PPTX slides."""
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError("Playwright is required for PPTX screenshot export.") from exc
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPTX export. Run: pip install python-pptx"
        ) from exc

    prs = Presentation()
    prs.slide_width = Emu(int(project.canvas_width * 914400 / 96))
    prs.slide_height = Emu(int(project.canvas_height * 914400 / 96))
    blank_layout = prs.slide_layouts[6]

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height},
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            png_bytes = pg.screenshot(full_page=False, type="png")
            ctx.close()

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), prs.slide_width, prs.slide_height)
            if page.notes:
                slide.notes_slide.notes_text_frame.text = page.notes
        browser.close()

    buf = io.BytesIO()
    prs.save(buf)
    out_path = describe_export_destination(project, "pptx", pages, directory=directory)
    return _save_bytes(out_path, buf.getvalue(), "PPTX")


def _css_length_to_px(value: Any, default: float = 0.0) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    match = _CSS_NUMBER_RE.search(str(value or ""))
    return float(match.group(0)) if match else default


def _px_to_emu(value: Any) -> int:
    return int(round(max(0.0, _css_length_to_px(value)) * 914400 / 96))


def _px_to_pt(value: Any, default: float = 12.0) -> float:
    pixels = _css_length_to_px(value, default / 0.75)
    return max(1.0, pixels * 72 / 96)


def _parse_css_color(value: Any):
    text = str(value or "").strip().lower()
    if not text or text == "transparent":
        return None
    if text.startswith("#"):
        hex_value = text[1:]
        if len(hex_value) == 3:
            hex_value = "".join(char * 2 for char in hex_value)
        if len(hex_value) == 6:
            return tuple(int(hex_value[i:i + 2], 16) for i in range(0, 6, 2)), 1.0
        return None
    channels = [float(number) for number in re.findall(r"[\d.]+", text)]
    if len(channels) < 3:
        return None
    rgb = tuple(max(0, min(255, int(round(channel)))) for channel in channels[:3])
    alpha = max(0.0, min(1.0, channels[3])) if len(channels) >= 4 else 1.0
    return rgb, alpha


def _font_name_from_css(value: Any) -> str | None:
    families = [family.strip().strip('"\'') for family in str(value or "").split(",")]
    return next((family for family in families if family), None)


def _weight_from_css(value: Any) -> int:
    text = str(value or "").strip().lower()
    if text == "bold":
        return 700
    if text == "normal":
        return 400
    match = _CSS_NUMBER_RE.search(text)
    return int(float(match.group(0))) if match else 400


def _alignment_from_css(value: Any):
    from pptx.enum.text import PP_ALIGN

    normalized = str(value or "left").strip().lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    if normalized == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def _apply_fill(fill, color_value: Any) -> bool:
    parsed = _parse_css_color(color_value)
    if not parsed:
        return False
    from pptx.dml.color import RGBColor

    rgb, alpha = parsed
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    try:
        fill.transparency = max(0.0, min(1.0, 1.0 - alpha))
    except Exception:
        pass
    return True


def _apply_line(shape, item: dict[str, Any]) -> None:
    parsed = _parse_css_color(item.get("borderColor"))
    border_width = _css_length_to_px(item.get("borderWidth"))
    if not parsed or border_width <= 0:
        shape.line.fill.background()
        return
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    rgb, _alpha = parsed
    shape.line.color.rgb = RGBColor(*rgb)
    shape.line.width = Pt(_px_to_pt(border_width))


# ---------------------------------------------------------------------------
# Native-shape editable PPTX helpers (Strategy A).
#
# These emit real DrawingML — gradient fills, translucent fills, outer
# shadows — so every card in the export is a first-class PowerPoint
# shape, not a picture.  Double-clicking a shape gives the Shape Tools
# ribbon, text inside cards is editable, and images keep Picture Tools.
# ---------------------------------------------------------------------------

_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DRAWINGML_NS_DECL = f' xmlns:a="{_DRAWINGML_NS}"'

_CSS_COLOR_TOKEN_RE = re.compile(
    r"#[0-9a-fA-F]{3,8}|rgba?\([^)]*\)|hsla?\([^)]*\)|[a-zA-Z][a-zA-Z0-9]*",
)
_LINEAR_GRADIENT_RE = re.compile(r"linear-gradient\((.*)\)\s*$", re.IGNORECASE | re.DOTALL)
_RADIAL_GRADIENT_RE = re.compile(r"radial-gradient\((.*)\)\s*$", re.IGNORECASE | re.DOTALL)


def _split_css_top_level(text: str, sep: str = ",") -> list[str]:
    """Split ``text`` on ``sep`` while ignoring separators inside parens."""
    parts: list[str] = []
    depth = 0
    buf: list[str] = []
    for ch in text:
        if ch == "(":
            depth += 1
            buf.append(ch)
        elif ch == ")":
            depth = max(0, depth - 1)
            buf.append(ch)
        elif ch == sep and depth == 0:
            parts.append("".join(buf).strip())
            buf = []
        else:
            buf.append(ch)
    if buf:
        parts.append("".join(buf).strip())
    return parts


def _parse_css_angle(value: str) -> float:
    """Return a CSS angle in degrees (0=up, increases clockwise)."""
    text = (value or "").strip().lower()
    if not text:
        return 180.0
    if text.startswith("to "):
        direction = text[3:].strip()
        mapping = {
            "top": 0.0, "right top": 45.0, "top right": 45.0,
            "right": 90.0, "bottom right": 135.0, "right bottom": 135.0,
            "bottom": 180.0, "bottom left": 225.0, "left bottom": 225.0,
            "left": 270.0, "top left": 315.0, "left top": 315.0,
        }
        return mapping.get(direction, 180.0)
    match = re.search(r"(-?\d+(?:\.\d+)?)", text)
    if not match:
        return 180.0
    degrees = float(match.group(1))
    if "rad" in text:
        degrees = degrees * 180.0 / math.pi
    elif "turn" in text:
        degrees = degrees * 360.0
    elif "grad" in text:
        degrees = degrees * 0.9
    return degrees


def _parse_gradient_stops(stops_text: list[str]) -> list[tuple[float, tuple[tuple[int, int, int], float]]]:
    """Parse CSS color stops into ``[(position_0_to_1, (rgb, alpha)), ...]``.

    ``getComputedStyle`` normalises every stop to ``rgb(...) NN%`` (or
    explicit pixel length), so the parser stays simple.  When a stop
    lacks an explicit position, it is interpolated evenly.
    """
    parsed: list[tuple[float | None, tuple[tuple[int, int, int], float]]] = []
    for raw in stops_text:
        token = raw.strip()
        if not token:
            continue
        # Find the color function (rgb/rgba/hsl/hsla/hex/keyword).
        color_match = re.match(
            r"(#[0-9a-fA-F]{3,8}|rgba?\([^)]*\)|hsla?\([^)]*\)|[a-zA-Z][a-zA-Z0-9]*)",
            token,
        )
        if not color_match:
            continue
        color_text = color_match.group(1)
        rest = token[color_match.end():].strip()
        parsed_color = _parse_css_color(color_text)
        if not parsed_color:
            continue
        position: float | None = None
        if rest:
            pos_match = re.search(r"(-?\d+(?:\.\d+)?)\s*(%|px)?", rest)
            if pos_match:
                raw_value = float(pos_match.group(1))
                unit = (pos_match.group(2) or "%").lower()
                if unit == "%":
                    position = raw_value / 100.0
                elif unit == "px":
                    # Without element size we cannot exactly resolve px
                    # stops; approximate 0 / middle / 1.  Good enough
                    # for the small handful of cases that use pixel
                    # stops.
                    position = None
        parsed.append((position, parsed_color))

    if not parsed:
        return []
    # Fill implicit positions evenly.
    out: list[tuple[float, tuple[tuple[int, int, int], float]]] = []
    n = len(parsed)
    for idx, (pos, color) in enumerate(parsed):
        if pos is None:
            pos = 0.0 if n == 1 else idx / (n - 1)
        pos = max(0.0, min(1.0, pos))
        out.append((pos, color))
    out.sort(key=lambda item: item[0])
    return out


def _parse_linear_gradient(value: str) -> dict[str, Any] | None:
    match = _LINEAR_GRADIENT_RE.search(value or "")
    if not match:
        return None
    parts = _split_css_top_level(match.group(1))
    if not parts:
        return None
    first = parts[0].strip().lower()
    if first.startswith("to ") or re.match(r"^-?\d", first):
        angle = _parse_css_angle(parts[0])
        stops_raw = parts[1:]
    else:
        angle = 180.0  # CSS default: top to bottom
        stops_raw = parts
    stops = _parse_gradient_stops(stops_raw)
    if len(stops) < 2:
        return None
    return {"type": "linear", "angle": angle, "stops": stops}


def _parse_radial_gradient(value: str) -> dict[str, Any] | None:
    match = _RADIAL_GRADIENT_RE.search(value or "")
    if not match:
        return None
    parts = _split_css_top_level(match.group(1))
    if not parts:
        return None
    # The first segment may be the shape/size/position prefix; treat
    # segments that contain a color token as stops, and skip the rest.
    stops_raw = [p for p in parts if re.search(r"#|rgb|hsl|[a-zA-Z]{3,}\s*[\d\-]", p)]
    if not stops_raw:
        stops_raw = parts
    stops = _parse_gradient_stops(stops_raw)
    if len(stops) < 2:
        return None
    return {"type": "radial", "stops": stops}


def _parse_box_shadow(value: str) -> list[dict[str, Any]]:
    """Parse a CSS ``box-shadow`` into a list of shadow dicts.

    We honour only outer drop shadows (inset is ignored because PowerPoint
    has no true inset shadow on auto-shapes; the effect would require
    multiple stacked shapes).  The first outer shadow is the most
    visually important one and the only one PowerPoint renders natively
    when multiple are specified.
    """
    text = (value or "").strip()
    if not text or text.lower() == "none":
        return []
    shadows = []
    for part in _split_css_top_level(text):
        segment = part.strip()
        if not segment or "inset" in segment.lower():
            continue
        # Extract color token (rgb/rgba/hsl/hex/keyword) first so the
        # remaining numbers are unambiguous.
        color_match = re.search(
            r"(#[0-9a-fA-F]{3,8}|rgba?\([^)]*\)|hsla?\([^)]*\))",
            segment,
        )
        color_token = color_match.group(1) if color_match else ""
        remainder = segment.replace(color_token, " ") if color_token else segment
        numbers = [
            float(m.group(0))
            for m in re.finditer(r"-?\d+(?:\.\d+)?", remainder)
        ]
        if len(numbers) < 2:
            continue
        x = numbers[0]
        y = numbers[1]
        blur = numbers[2] if len(numbers) >= 3 else 0.0
        spread = numbers[3] if len(numbers) >= 4 else 0.0
        parsed_color = _parse_css_color(color_token) if color_token else ((0, 0, 0), 0.6)
        if not parsed_color:
            parsed_color = ((0, 0, 0), 0.6)
        rgb, alpha = parsed_color
        shadows.append({
            "x": x, "y": y, "blur": blur, "spread": spread,
            "rgb": rgb, "alpha": alpha,
        })
    return shadows


def _hex_rgb(rgb: tuple[int, int, int]) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _build_solid_fill_xml(rgb: tuple[int, int, int], alpha: float) -> str:
    hex_rgb = _hex_rgb(rgb)
    alpha_clamped = max(0.0, min(1.0, alpha))
    if alpha_clamped >= 0.999:
        return f'<a:solidFill{_DRAWINGML_NS_DECL}><a:srgbClr val="{hex_rgb}"/></a:solidFill>'
    alpha_val = int(round(alpha_clamped * 100000))
    return (
        f'<a:solidFill{_DRAWINGML_NS_DECL}>'
        f'<a:srgbClr val="{hex_rgb}"><a:alpha val="{alpha_val}"/></a:srgbClr>'
        f'</a:solidFill>'
    )


def _build_gradient_stop_xml(stop: tuple[float, tuple[tuple[int, int, int], float]]) -> str:
    pos, (rgb, alpha) = stop
    pos_val = int(round(max(0.0, min(1.0, pos)) * 100000))
    hex_rgb = _hex_rgb(rgb)
    if alpha >= 0.999:
        inner = f'<a:srgbClr val="{hex_rgb}"/>'
    else:
        alpha_val = int(round(max(0.0, min(1.0, alpha)) * 100000))
        inner = (
            f'<a:srgbClr val="{hex_rgb}">'
            f'<a:alpha val="{alpha_val}"/>'
            f'</a:srgbClr>'
        )
    return f'<a:gs pos="{pos_val}">{inner}</a:gs>'


def _build_linear_gradient_xml(spec: dict[str, Any]) -> str:
    # CSS angle: 0deg = up, clockwise.  OOXML <a:lin ang="...">: 0 =
    # east, clockwise, unit 60000ths of a degree.  Mapping:
    # ooxml = (css - 90) mod 360.
    css_angle = float(spec.get("angle") or 180.0)
    ooxml_angle = int(round(((css_angle - 90.0) % 360.0) * 60000))
    stops_xml = "".join(_build_gradient_stop_xml(s) for s in spec["stops"])
    return (
        f'<a:gradFill{_DRAWINGML_NS_DECL} flip="none" rotWithShape="1">'
        f'<a:gsLst>{stops_xml}</a:gsLst>'
        f'<a:lin ang="{ooxml_angle}" scaled="0"/>'
        f'</a:gradFill>'
    )


def _build_radial_gradient_xml(spec: dict[str, Any]) -> str:
    stops_xml = "".join(_build_gradient_stop_xml(s) for s in spec["stops"])
    return (
        f'<a:gradFill{_DRAWINGML_NS_DECL} flip="none" rotWithShape="1">'
        f'<a:gsLst>{stops_xml}</a:gsLst>'
        f'<a:path path="circle">'
        f'<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>'
        f'</a:path>'
        f'</a:gradFill>'
    )


def _build_outer_shadow_xml(shadow: dict[str, Any]) -> str:
    x = float(shadow.get("x") or 0.0)
    y = float(shadow.get("y") or 0.0)
    blur = float(shadow.get("blur") or 0.0)
    rgb = shadow.get("rgb") or (0, 0, 0)
    alpha = float(shadow.get("alpha") or 1.0)
    dist_px = math.hypot(x, y)
    if dist_px > 0:
        dir_deg = math.degrees(math.atan2(y, x)) % 360.0
    else:
        dir_deg = 90.0  # down
    dir_val = int(round(dir_deg * 60000))
    dist_emu = int(round(dist_px * 914400 / 96))
    blur_emu = int(round(blur * 914400 / 96))
    hex_rgb = _hex_rgb(rgb)
    alpha_val = int(round(max(0.0, min(1.0, alpha)) * 100000))
    return (
        f'<a:effectLst{_DRAWINGML_NS_DECL}>'
        f'<a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="{dir_val}" '
        f'algn="ctr" rotWithShape="0">'
        f'<a:srgbClr val="{hex_rgb}"><a:alpha val="{alpha_val}"/></a:srgbClr>'
        f'</a:outerShdw>'
        f'</a:effectLst>'
    )


_FILL_TAGS = ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill")


def _replace_shape_fill_xml(shape, fill_xml: str) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    sp_pr = shape.fill._xPr  # spPr
    for tag in _FILL_TAGS:
        for el in sp_pr.findall(qn(f"a:{tag}")):
            sp_pr.remove(el)
    new_el = etree.fromstring(fill_xml)
    # Insert fill before <a:ln> / <a:effectLst> / <a:scene3d>, after geometry.
    anchor_index = len(sp_pr)
    for idx, child in enumerate(sp_pr):
        if child.tag in (qn("a:ln"), qn("a:effectLst"), qn("a:scene3d"), qn("a:sp3d")):
            anchor_index = idx
            break
    sp_pr.insert(anchor_index, new_el)


def _apply_effect_xml(shape, effect_xml: str) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    sp_pr = shape.fill._xPr  # spPr
    # Remove any existing effectLst so we don't stack.
    for el in sp_pr.findall(qn("a:effectLst")):
        sp_pr.remove(el)
    new_el = etree.fromstring(effect_xml)
    # effectLst must come after fill and line.
    insertion_index = len(sp_pr)
    for idx, child in enumerate(sp_pr):
        if child.tag in (qn("a:scene3d"), qn("a:sp3d")):
            insertion_index = idx
            break
    sp_pr.insert(insertion_index, new_el)


def _apply_native_fill(shape, item: dict[str, Any]) -> bool:
    """Apply a native DrawingML fill matching the CSS rendering.

    Returns True if any fill was applied.  When the element has no
    meaningful visual fill (e.g. transparent bg + no gradient), calls
    ``shape.fill.background()`` so PowerPoint shows a hollow shape.
    """
    bg_image = str(item.get("backgroundImage") or "").strip()
    if bg_image and bg_image.lower() != "none":
        gradient = _parse_linear_gradient(bg_image) or _parse_radial_gradient(bg_image)
        if gradient:
            xml = (
                _build_linear_gradient_xml(gradient)
                if gradient["type"] == "linear"
                else _build_radial_gradient_xml(gradient)
            )
            _replace_shape_fill_xml(shape, xml)
            return True
    bg_color = item.get("backgroundColor")
    parsed = _parse_css_color(bg_color)
    if parsed and parsed[1] > 0.02:
        rgb, alpha = parsed
        _replace_shape_fill_xml(shape, _build_solid_fill_xml(rgb, alpha))
        return True
    shape.fill.background()
    return False


def _apply_native_shadow(shape, item: dict[str, Any]) -> None:
    shadows = _parse_box_shadow(str(item.get("boxShadow") or ""))
    if not shadows:
        return
    # PowerPoint only renders the first outer shadow on an auto-shape.
    _apply_effect_xml(shape, _build_outer_shadow_xml(shadows[0]))



def _shape_type_for_radius(item: dict[str, Any]):
    from pptx.enum.shapes import MSO_SHAPE

    radius = _css_length_to_px(item.get("borderRadius"))
    return MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0.5 else MSO_SHAPE.RECTANGLE


def _text_vertical_anchor(item: dict[str, Any]):
    from pptx.enum.text import MSO_VERTICAL_ANCHOR

    display = str(item.get("display") or "").strip().lower()
    align_items = str(item.get("alignItems") or "").strip().lower()
    if display == "flex" and align_items in {"center", "flex-end"}:
        return MSO_VERTICAL_ANCHOR.MIDDLE if align_items == "center" else MSO_VERTICAL_ANCHOR.BOTTOM
    return MSO_VERTICAL_ANCHOR.TOP


def _collect_native_slide(project: DesignerProject, page_html: str, *, page_index: int):
    """Collect DOM items + per-element raster screenshots for <img>/<svg>.

    Native mode: no full-slide raster, no per-card snapshots.  Every
    visible element becomes a native PPTX shape (for fills/borders/
    gradients/shadows) OR an editable textbox (for text) OR — only for
    ``<img>`` and inline ``<svg>`` — a localized screenshot.  Those two
    tags have no visible HTML descendants that would also emit, so
    there is no ancestor/descendant duplication possible.
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError(
            "Playwright is required for editable PPTX export. Run: pip install playwright && playwright install chromium"
        ) from exc

    html = render_page_html(project, page_html, page_index=page_index)
    screenshot_bytes: dict[str, bytes] = {}

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": project.canvas_width, "height": project.canvas_height},
            device_scale_factor=2,
        )
        pg = ctx.new_page()
        pg.set_content(html, wait_until="networkidle")
        snapshot = pg.evaluate(_RENDERED_DOM_EXPORT_SCRIPT)

        # Render every raster item in an ISOLATED page.  Using
        # locator.screenshot() on the live page would clip the page
        # at the element rect, which captures any overlapping text or
        # shapes rendered on top of an <img>/<svg> — causing visible
        # duplication in the final PPTX.  Isolated rendering
        # guarantees the raster contains only the element itself.
        for item in snapshot.get("items", []):
            if item.get("kind") != "raster":
                continue
            screenshot_id = item.get("screenshotId")
            if not screenshot_id:
                continue
            width = max(1, int(round(float(item.get("width") or 1))))
            height = max(1, int(round(float(item.get("height") or 1))))
            opacity = float(item.get("opacity") or 1.0)
            try:
                if item.get("svgOuterHTML"):
                    body = (
                        f'<div style="width:{width}px;height:{height}px;opacity:{opacity}">'
                        f'{item["svgOuterHTML"]}</div>'
                    )
                else:
                    src = str(item.get("src") or "")
                    if not src:
                        continue
                    fit = str(item.get("objectFit") or "fill")
                    pos = str(item.get("objectPosition") or "50% 50%")
                    src_attr = src.replace('"', '&quot;')
                    body = (
                        f'<img src="{src_attr}" '
                        f'style="width:{width}px;height:{height}px;'
                        f'object-fit:{fit};object-position:{pos};'
                        f'opacity:{opacity};display:block">'
                    )
                iso_html = (
                    '<!DOCTYPE html><html><head><meta charset="utf-8">'
                    '<style>html,body{margin:0;padding:0;background:transparent;}</style>'
                    f'</head><body>{body}</body></html>'
                )
                iso = ctx.new_page()
                iso.set_viewport_size({"width": width, "height": height})
                iso.set_content(iso_html, wait_until="networkidle")
                screenshot_bytes[screenshot_id] = iso.screenshot(
                    full_page=False, omit_background=True, clip={
                        "x": 0, "y": 0, "width": width, "height": height,
                    },
                )
                iso.close()
            except Exception:
                logger.debug("Isolated raster render failed for %s", screenshot_id, exc_info=True)

        ctx.close()
        browser.close()

    return snapshot, screenshot_bytes


def _rendered_item_sort_key(item: dict[str, Any]) -> tuple[int, int, int]:
    # Shapes render first (behind), then rasters, then text on top.
    kind_order = {"shape": 0, "raster": 1, "text": 2}.get(item.get("kind") or "", 3)
    return int(item.get("zIndex") or 0), int(item.get("order") or 0), kind_order


def _dedupe_text_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Collapse duplicate text emissions.

    Even with the JS-side inline-ownership logic, edge cases slip
    through (e.g. decorative outline-text layered on the real heading,
    an accent span that is a substring of the full heading, or a
    ``::before`` content strip duplicated by an overlay element).
    We run a final safety net on the Python side that catches two
    patterns:

    1. **Exact text + rect overlap** — two items with the same
       normalised text whose bounding rects overlap by >30%% of the
       smaller area.  Classic stroke-over-fill duplication effect.
    2. **Containment** — one text is a substring of the other AND
       their rects overlap by >50%% of the smaller area.  Happens when
       the agent wraps part of a heading in a coloured ``<span>`` that
       paints with ``display:block`` (so the JS dedup classifies it as
       a block child and the parent switches to own-text mode) AND the
       parent's own text is left carrying the full string, producing
       one full-string emission plus one accent-fragment emission.

    In both cases we keep the item that appeared first in DOM order
    (typically the outer / larger element) and drop the duplicate.
    """
    def _rect(item):
        try:
            return (
                float(item.get("x") or 0.0),
                float(item.get("y") or 0.0),
                float(item.get("width") or 0.0),
                float(item.get("height") or 0.0),
            )
        except (TypeError, ValueError):
            return None

    def _norm(text):
        return " ".join(str(text or "").split())

    text_positions: list[int] = []  # indices into `items`
    for idx, it in enumerate(items):
        if it.get("kind") == "text" and _norm(it.get("text")):
            text_positions.append(idx)

    drop: set[int] = set()
    for a_pos_i, a_idx in enumerate(text_positions):
        if a_idx in drop:
            continue
        a = items[a_idx]
        a_text = _norm(a.get("text"))
        a_rect = _rect(a)
        if a_rect is None:
            continue
        ax, ay, aw, ah = a_rect
        a_area = max(1.0, aw * ah)
        a_lower = a_text.lower()
        for b_idx in text_positions[a_pos_i + 1:]:
            if b_idx in drop:
                continue
            b = items[b_idx]
            b_text = _norm(b.get("text"))
            b_rect = _rect(b)
            if b_rect is None:
                continue
            bx, by, bw, bh = b_rect
            b_area = max(1.0, bw * bh)
            b_lower = b_text.lower()
            ox = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
            oy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
            if ox <= 0 or oy <= 0:
                continue
            overlap_area = ox * oy
            smaller = min(a_area, b_area)
            overlap_ratio = overlap_area / smaller
            if a_lower == b_lower and overlap_ratio > 0.3:
                drop.add(b_idx)
                continue
            if overlap_ratio > 0.5 and a_lower and b_lower and a_lower != b_lower:
                # Substring containment — the shorter text is a
                # fragment of the longer.  Drop the shorter so the
                # full heading survives without its accent fragment
                # rendered twice on top.
                if a_lower in b_lower:
                    drop.add(a_idx)
                    break
                if b_lower in a_lower:
                    drop.add(b_idx)
                    continue

    return [it for idx, it in enumerate(items) if idx not in drop]


def _add_rendered_item_to_slide(slide, item: dict[str, Any], screenshots: dict[str, bytes]) -> None:
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_AUTO_SIZE

    left = Emu(_px_to_emu(item.get("x")))
    top = Emu(_px_to_emu(item.get("y")))
    width = Emu(_px_to_emu(item.get("width")))
    height = Emu(_px_to_emu(item.get("height")))
    kind = item.get("kind")

    if kind == "raster":
        screenshot_id = item.get("screenshotId")
        png_bytes = screenshots.get(screenshot_id or "")
        if png_bytes:
            slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)
        return

    if kind == "shape":
        shape = slide.shapes.add_shape(_shape_type_for_radius(item), left, top, width, height)
        _apply_native_fill(shape, item)
        _apply_line(shape, item)
        _apply_native_shadow(shape, item)
        return

    if kind != "text":
        return

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = str(item.get("whiteSpace") or "normal").strip().lower() not in {"nowrap", "pre"}
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Emu(_px_to_emu(item.get("paddingLeft")))
    text_frame.margin_right = Emu(_px_to_emu(item.get("paddingRight")))
    text_frame.margin_top = Emu(_px_to_emu(item.get("paddingTop")))
    text_frame.margin_bottom = Emu(_px_to_emu(item.get("paddingBottom")))
    text_frame.vertical_anchor = _text_vertical_anchor(item)

    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(item.get("text") or "")
    paragraph.alignment = _alignment_from_css(item.get("textAlign"))
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)

    line_height = str(item.get("lineHeight") or "").strip().lower()
    if line_height and line_height != "normal":
        paragraph.line_spacing = Pt(_px_to_pt(line_height))

    font = paragraph.font
    font.name = _font_name_from_css(item.get("fontFamily"))
    font.size = Pt(_px_to_pt(item.get("fontSize"), default=14.0))
    font.bold = _weight_from_css(item.get("fontWeight")) >= 600
    font.italic = str(item.get("fontStyle") or "").strip().lower() == "italic"

    parsed_color = _parse_css_color(item.get("color"))
    if parsed_color:
        from pptx.dml.color import RGBColor

        rgb, _alpha = parsed_color
        font.color.rgb = RGBColor(*rgb)


def export_pptx_structured(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render the Designer preview as a fully-editable native PPTX.

    Every card, pill, badge, and divider becomes a real PowerPoint
    auto-shape (with native gradient / solid / translucent fills, border,
    and drop shadow); every text run becomes an editable textbox; only
    ``<img>`` bitmaps and inline ``<svg>`` icons are embedded as
    pictures.  Double-clicking any card in PowerPoint opens the Shape
    Tools ribbon — colors, sizes, and copy are all user-editable.

    For pixel-perfect visual parity use the High-Fidelity mode instead
    (``export_pptx_screenshot``), which embeds full-slide screenshots.
    """
    selected = _selected_pages(project, pages)

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPTX export. Run: pip install python-pptx"
        ) from exc
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(int(project.canvas_width * 914400 / 96))
    prs.slide_height = Emu(int(project.canvas_height * 914400 / 96))
    blank_layout = prs.slide_layouts[6]

    for page_index, page in selected:
        slide = prs.slides.add_slide(blank_layout)

        snapshot, screenshots = _collect_native_slide(
            project, page.html, page_index=page_index,
        )

        # Slide background: apply solid body color when present.  CSS
        # gradients on the body/root (rare — the Designer paints its
        # background on a full-slide card div, which becomes a native
        # shape further down) fall through to the default slide bg.
        _apply_fill(slide.background.fill, snapshot.get("backgroundColor"))

        items = sorted(snapshot.get("items", []), key=_rendered_item_sort_key)
        items = _dedupe_text_items(items)
        for item in items:
            try:
                _add_rendered_item_to_slide(slide, item, screenshots)
            except Exception:
                logger.debug("Failed to add item %s", item.get("kind"), exc_info=True)

        if page.notes:
            slide.notes_slide.notes_text_frame.text = page.notes

    buf = io.BytesIO()
    prs.save(buf)
    out_path = describe_export_destination(project, "pptx", pages, mode="structured", directory=directory)
    return _save_bytes(out_path, buf.getvalue(), "Editable PPTX")